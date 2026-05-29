"""
本地文件解析模組：將上傳的檔案轉換為可切分的文字。
PDF 一律經 opendataloader-pdf（fast / Java）輸出 Markdown，再去做噪。
需本機 Java 11+。
"""
import re
import tempfile
from pathlib import Path

# PDF 垂直水印常拆成單字元行（如 Notion/BuildMoat 匯出的 grow.tao 側欄）
_PDF_GLYPH_NOISE_LINE = re.compile(r"^[\.\-a-z]{1,2}$")
# 側欄水印整行
_PDF_WATERMARK_LINE = re.compile(r"^(?:moat\.org|buildmoat\.org|biuld)\s*$", re.I)
# 僅圖片引用、無文字描述的 markdown 行（不做 OCR / 圖片描述）
_PDF_IMAGE_ONLY_LINE = re.compile(r"^\s*!\[[^\]]*\]\(<[^>]+>\)\s*$")
# 側欄水印插入中文詞內：不l均、分散i式、加入時u會、簡單m。
_PDF_CJK_INLINE_ASCII = re.compile(r"([\u4e00-\u9fff])([a-z])([\u4e00-\u9fff])")
# 中文後多一個 ascii（簡單m、加入時u）再接標點/空白/中文
_PDF_CJK_TRAILING_ASCII = re.compile(
    r"([\u4e00-\u9fff])([a-z])(?=[\u4e00-\u9fff。，、；：！？\.\s\n]|$)"
)
# 可能是.幾百萬
_PDF_CJK_DOT_CJK = re.compile(r"([\u4e00-\u9fff])\.([\u4e00-\u9fff])")
# 規則很簡單： l
_PDF_COLON_STRAY_ASCII = re.compile(r"([：:])\s*[a-z](?=\s*(?:\n|[1-9]|$))")
# 英文詞後多餘字母：cache ring t、16384 t
_PDF_EN_TRAILING_T = re.compile(r"(\w)\s+t(?=\s*(?:\n|->|$))")
# Cluster r就是
_PDF_EN_R_BEFORE_CJK = re.compile(r"([A-Za-z]{2,})\s+r([\u4e00-\u9fff])")
# 行首 orphan l user_id
_PDF_LINE_LEADING_ASCII = re.compile(r"(?m)^[a-z]\s+(?=[a-z_])")
_PDF_INLINE_FIXES: tuple[tuple[re.Pattern[str], str], ...] = (
    (re.compile(r"hash\(ukey\)", re.I), "hash(key)"),
    (re.compile(r"\bhashd\b", re.I), "hash"),
    (re.compile(r"\bhbash\b", re.I), "hash"),
    (re.compile(r"\bhashinog\b", re.I), "hashing"),
    (re.compile(r"\bmodurlo\b", re.I), "modulo"),
    (re.compile(r"\bhotu\b", re.I), "hot"),
    (re.compile(r"grow\.tao", re.I), ""),
    (re.compile(r"\bcaler\b", re.I), "caller"),
    (re.compile(r"\bwebhok\b", re.I), "webhook"),
    (re.compile(r"\bHTP\b"), "HTTP"),
    (re.compile(r"\btoling\b", re.I), "tooling"),
    (re.compile(r"\bcomand\b", re.I), "command"),
    (re.compile(r"\beventualy\b", re.I), "eventually"),
    (re.compile(r"\beror\b", re.I), "error"),
    (re.compile(r"\bofset\b", re.I), "offset"),
    (re.compile(r"\bchanel\b", re.I), "channel"),
)


def extract_text(filename: str, raw_bytes: bytes) -> str:
    """
    從檔案 bytes 抽取文字。
    優先使用對應格式的 parser；若失敗，fallback 到 utf-8 decode。
    """
    suffix = Path(filename).suffix.lower()
    try:
        if suffix in (".txt", ".md"):
            return _extract_text_plain(raw_bytes)
        elif suffix == ".pdf":
            return _extract_pdf(raw_bytes, filename=filename)
        elif suffix == ".docx":
            return _extract_docx(raw_bytes)
        elif suffix == ".pptx":
            return _extract_pptx(raw_bytes)
        elif suffix in (".html", ".htm"):
            return _extract_html(raw_bytes)
        elif suffix == ".epub":
            # EPUB 在 upload router 已被切成多章節獨立檔；這條路徑只剩本機 chunker
            # 測試 / 一次性脫離 upload pipeline 的工具會走，因此用 epub_splitter
            # 並把章節接起來給 chunker。
            from .epub_splitter import split_epub_by_toc
            return "\n\n".join(text for _, text in split_epub_by_toc(raw_bytes))
        else:
            return _fallback_decode(raw_bytes)
    except Exception:
        return _fallback_decode(raw_bytes)


def _extract_text_plain(raw: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "big5", "gb18030"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_pdf(raw: bytes, *, filename: str) -> str:
    """opendataloader-pdf fast 模式 → Markdown（不啟用 hybrid / OCR）。"""
    import opendataloader_pdf

    safe_name = Path(filename).name or "upload.pdf"
    if not safe_name.lower().endswith(".pdf"):
        safe_name = f"{safe_name}.pdf"

    with tempfile.TemporaryDirectory(prefix="odpdf_") as tmpdir:
        tmp = Path(tmpdir)
        pdf_path = tmp / safe_name
        pdf_path.write_bytes(raw)
        out_dir = tmp / "out"
        out_dir.mkdir()
        opendataloader_pdf.convert(
            input_path=[str(pdf_path)],
            output_dir=str(out_dir),
            format="markdown",
        )
        md_files = sorted(out_dir.glob("*.md"))
        if not md_files:
            raise RuntimeError("opendataloader-pdf 未產生 markdown 輸出")
        text = md_files[0].read_text(encoding="utf-8", errors="replace")
    return postprocess_pdf_markdown(text)


def postprocess_pdf_markdown(text: str) -> str:
    """PDF Markdown 去噪：水印行、圖片佔位、行內水印殘留、常見 OCR 誤字。"""
    kept: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            kept.append("")
            continue
        if _PDF_WATERMARK_LINE.match(stripped):
            continue
        if _PDF_GLYPH_NOISE_LINE.match(stripped):
            continue
        if _PDF_IMAGE_ONLY_LINE.match(stripped):
            continue
        kept.append(line.rstrip())

    merged = "\n".join(kept)
    merged = re.sub(r"\n{4,}", "\n\n\n", merged)
    merged = _fix_inline_watermark_chars(merged)

    for pattern, replacement in _PDF_INLINE_FIXES:
        merged = pattern.sub(replacement, merged)
    merged = re.sub(r"  +", " ", merged)
    return merged.strip()


# 測試與舊呼叫點相容
_clean_pdf_text = postprocess_pdf_markdown


def _fix_inline_watermark_chars(text: str) -> str:
    """移除側欄水印插入中文/英文詞內的孤立 ascii。"""
    prev = None
    while prev != text:
        prev = text
        text = _PDF_CJK_INLINE_ASCII.sub(r"\1\3", text)
        text = _PDF_CJK_TRAILING_ASCII.sub(r"\1", text)
    text = _PDF_CJK_DOT_CJK.sub(r"\1\2", text)
    text = _PDF_COLON_STRAY_ASCII.sub(r"\1", text)
    text = _PDF_EN_TRAILING_T.sub(r"\1", text)
    text = _PDF_EN_R_BEFORE_CJK.sub(r"\1 \2", text)
    text = _PDF_LINE_LEADING_ASCII.sub("", text)
    return text


def _extract_docx(raw: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(raw))
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        # 保留標題層級
        if para.style and para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading ", "").strip()
            try:
                hashes = "#" * int(level)
            except ValueError:
                hashes = "##"
            parts.append(f"{hashes} {text}")
        else:
            parts.append(text)
    return "\n\n".join(parts)


def _extract_pptx(raw: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_html(raw: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(raw, "lxml")
    for tag in soup(["script", "style", "head", "nav", "footer"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _fallback_decode(raw: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "big5", "gb18030", "latin-1"):
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")
